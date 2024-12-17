import textwrap
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from helpers_af.merge_cells_if_needed import merge_cells_if_needed


def wrap_text_first_row(text):
    return '\n'.join(text.split())


def wrap_text(text, width):
    return '\n'.join(textwrap.wrap(text, width))



def merge_day_cells(table, day_values, time_values):
    day_count = len(day_values)
    col = 0
    while col < day_count:
        current_day = day_values[col]
        current_time = time_values[col]
        
        if current_time == "AM":
            if col + 1 < day_count and day_values[col + 1] == current_day:
                merge_cells_if_needed(table, 0, col + 1, 0, col + 2)
                cell = table.cell(0, col + 1)
                cell.text = current_day if current_day is not None else ""
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(8)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                col += 2
            else:
                cell = table.cell(0, col + 1)
                cell.text = current_day if current_day is not None else ""
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(8)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                col += 1
        else:
            cell = table.cell(0, col + 1)
            cell.text = current_day if current_day is not None else ""
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            col += 1