from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from helpers_af.get_af_atmospherics import get_af_atmospherics
from helpers_af.get_af_temperature import get_af_temperature
from helpers_af.create_af_forecast_df import create_af_forecast_df
from helpers_af.add_dataframe_to_slide import add_dataframe_to_slide_single_city
from helpers_af.get_af_precipitation import get_af_precipitation
from helpers_af.add_emoji_to_table import add_emoji_to_table
from helpers_af.create_af_map import create_map
from helpers_af.df_for_slide import df_for_slide

def create_powerpoint_single_city(
        location_name,
        location1_html_from_af,
        location_illum_df,
        location_vis_condition,
        location_latitude, 
        location_longitude,
        start_day,
        num_days,
        zoom_level
    ):
    prs = Presentation()

    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Weather Forecast"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.underline = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    title_box.line.color.rgb = RGBColor(0, 0, 0)
    title_box.line.width = Pt(2)

    def add_textbox(slide, left, top, width, height, location_name, forecast_data, location_illum_df, location_vis_condition):
        # Pass both start_day and num_days to get_af_precipitation
        precipitation = get_af_precipitation(location1_html_from_af, start_day, num_days)
        temperature = get_af_temperature(location_forecast_df, num_days)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        tf.clear()
        p = tf.paragraphs[0]
        p.text = location_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.underline = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

        def add_second_level_title(tf, text):
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(14)
            p.font.bold = False
            p.font.underline = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            return p

        def add_second_level_data(tf, data):
            p = tf.add_paragraph()
            p.text = f"•   {data}"
            p.font.size = Pt(14)
            p.font.bold = False
            p.font.underline = False
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            return p
            
        add_second_level_title(tf, 'Precipitation')
        add_second_level_data(tf, precipitation)
        add_second_level_title(tf, 'Temperature')
        add_second_level_data(tf, temperature)
        add_second_level_title(tf, 'Atmospheric Conditions')

        atmospherics = get_af_atmospherics(forecast_data, location_illum_df, location_vis_condition)
        for key, value in atmospherics.items():
            if value:
                add_second_level_data(tf, f"{key}: {value}")
            else:
                add_second_level_data(tf, key)

        textbox.line.color.rgb = RGBColor(0, 0, 0)
        textbox.line.width = Pt(2)

    # Determine the width of the forecast table
    forecast_table_width = Inches(6.64)

    location_forecast_df = create_af_forecast_df(location1_html_from_af)
    add_textbox(slide, Inches(0), Inches(0.8), forecast_table_width, Inches(3.0), location_name, location_forecast_df, location_illum_df, location_vis_condition)

   
    slide_df = df_for_slide(location_forecast_df)
    table1, table1_illum = add_dataframe_to_slide_single_city(slide_df, location_illum_df, slide)

    for col_idx, emoji in enumerate(slide_df.loc['Graphic']):
        add_emoji_to_table(emoji, slide, table1, 2, col_idx + 1)

    def add_combined_border(slide, table1, table1_illum, border_width, outline_color):
        x1, y1, cx1, cy1 = table1._graphic_frame.left, table1._graphic_frame.top, table1._graphic_frame.width, table1._graphic_frame.height
        x2, y2, cx2, cy2 = table1_illum._graphic_frame.left, table1_illum._graphic_frame.top, table1_illum._graphic_frame.width, table1_illum._graphic_frame.height

        x = min(x1, x2)
        y = min(y1, y2)
        cx = max(x1 + cx1, x2 + cx2) - x
        cy = max(y1 + cy1, y2 + cy2) - y

        border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
        border_shape.fill.background()
        border_shape.line.color.rgb = outline_color
        border_shape.line.width = Pt(border_width)

    add_combined_border(slide, table1, table1_illum, 2, RGBColor(0, 0, 0))

    def add_red_box_over_table(slide, table, start_day, num_days):
        rows = table.rows
        cols = table.columns
        day_row_idx = 0
        gusts_row_idx = len(rows) - 1
        num_actual_cols = len(cols) - 1  # Exclude the header column

        # Ensure num_days is an integer and adjust to fit the table if needed
        num_days = int(num_days) * 2
        start_day = int(start_day) * 2 - 1  # Adjust to zero-based index
        
        if num_days + start_day > num_actual_cols:
            num_days = num_actual_cols - start_day

        # Calculate the left position and width
        x = table._graphic_frame.left + sum([table.columns[i + 1].width for i in range(start_day)])
        cx = sum([table.columns[i + 1].width for i in range(start_day, start_day + num_days)])

        # Calculate the top position and height
        y = table._graphic_frame.top
        cy = sum([table.rows[i].height for i in range(day_row_idx, gusts_row_idx + 1)])

        border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
        border_shape.fill.background()
        border_shape.line.color.rgb = RGBColor(255, 0, 0)
        border_shape.line.width = Pt(1)

    # Add red box over the table depending on operation length
    add_red_box_over_table(slide, table1, start_day, num_days)

    # Position the map to the right of the forecast table
    map_left = table1._graphic_frame.left + table1._graphic_frame.width
    map_top = Inches(0.8)
    map_width = prs.slide_width - map_left
    map_height = prs.slide_height - map_top

    map_image_path = create_map(location_latitude, location_longitude, zoom=zoom_level)
    map_img = slide.shapes.add_picture(map_image_path, map_left, map_top, width=map_width, height=map_height)

    # Add a black border around the map
    border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, map_left, map_top, map_width, map_height)
    border_shape.fill.background()
    border_shape.line.color.rgb = RGBColor(0, 0, 0)
    border_shape.line.width = Pt(2)

    slide.shapes._spTree.remove(map_img._element)
    slide.shapes._spTree.insert(2, map_img._element)

    prs.save(".output/Weather_Forecast.pptx")
