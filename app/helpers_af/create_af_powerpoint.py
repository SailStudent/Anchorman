from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from helpers_af.get_af_atmospherics import get_af_atmospherics
from helpers_af.get_af_temperature import get_af_temperature
from helpers_af.create_af_forecast_df import create_af_forecast_df
from helpers_af.add_af_dataframe_to_slide import add_af_dataframe_to_slide
from helpers_af.get_af_precipitation import get_af_precipitation
from helpers_af.add_emoji_to_table import add_emoji_to_table
from helpers_af.create_af_map import create_map_two_cities
from helpers_af.df_for_slide import df_for_slide

def create_af_powerpoint(
        prs,
        location1_name,
        location2_name,
        location1_html_from_af,
        location2_html_from_af,
        location1_illum_df,
        location2_illum_df,
        location1_vis_condition,
        location2_vis_condition,
        location1_lat, location1_lon,
        location2_lat, location2_lon,
        start_day,
        num_days
    ):
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = Inches(0.8)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "Weather Forecast"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.underline = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    title_box.line.color.rgb = RGBColor(0, 0, 0)
    title_box.line.width = Pt(2)

    def add_textbox(slide, left, top, width, height, location_name, forecast_data, illum_df, vis_condition, outline_color, html_file):
        precipitation = get_af_precipitation(html_file, start_day, num_days)  # Pass start_day and num_days
        temperature = get_af_temperature(forecast_data, num_days)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame

        tf.clear()
        p = tf.paragraphs[0]
        p.text = location_name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.underline = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = PP_ALIGN.CENTER

        def add_second_level_title(tf, text):
            p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(14)
            p.font.bold = False
            p.font.underline = True
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            return p

        def add_second_level_data(tf, data):
            p = tf.add_paragraph()
            p.text = f"•   {data}"
            p.font.size = Pt(14)
            p.font.bold = False
            p.font.underline = False
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.alignment = PP_ALIGN.LEFT
            return p
            
        add_second_level_title(tf, 'Precipitation')
        add_second_level_data(tf, precipitation)
        add_second_level_title(tf, 'Temperature')
        add_second_level_data(tf, temperature)
        add_second_level_title(tf, 'Atmospheric Conditions')

        atmospherics = get_af_atmospherics(forecast_data, illum_df, vis_condition)
        for key, value in atmospherics.items():
            if value:
                add_second_level_data(tf, f"{key}: {value}")
            else:
                add_second_level_data(tf, key)

        textbox.line.color.rgb = outline_color
        textbox.line.width = Pt(2)

    location1_forecast_df = create_af_forecast_df(location1_html_from_af)
    location2_forecast_df = create_af_forecast_df(location2_html_from_af)

    add_textbox(slide, Inches(0), Inches(0.8), Inches(3.695), Inches(3.0), location1_name, location1_forecast_df, location1_illum_df, location1_vis_condition, RGBColor(0, 128, 0), location1_html_from_af)
    add_textbox(slide, Inches(9.63), Inches(0.8), Inches(3.7), Inches(3.0), location2_name, location2_forecast_df, location2_illum_df, location2_vis_condition, RGBColor(255, 215, 0), location2_html_from_af)

    # Map creation for two cities
    map_image_path = create_map_two_cities(location1_lat, location1_lon, location2_lat, location2_lon)
    map_img = slide.shapes.add_picture(map_image_path, Inches(3.7), Inches(0.8), width=Inches(5.92), height=Inches(3.0))

    slide.shapes._spTree.remove(map_img._element)
    slide.shapes._spTree.insert(2, map_img._element)

    slide_loc1_df = df_for_slide(location1_forecast_df)
    slide_loc2_df = df_for_slide(location2_forecast_df)
    

    table1, table2, table1_illum, table2_illum = add_af_dataframe_to_slide(prs, slide, slide_loc1_df, slide_loc2_df, location1_illum_df, location2_illum_df)

    if 'Graphic' in location1_forecast_df.index:
        for col_idx, emoji in enumerate(location1_forecast_df.loc['Graphic']):
            add_emoji_to_table(emoji, slide, table1, 2, col_idx + 1)

    if 'Graphic' in location2_forecast_df.index:
        for col_idx, emoji in enumerate(location2_forecast_df.loc['Graphic']):
            add_emoji_to_table(emoji, slide, table2, 2, col_idx + 1)

    def add_combined_border(slide, table1, table2, border_width, outline_color):
        x1, y1, cx1, cy1 = table1._graphic_frame.left, table1._graphic_frame.top, table1._graphic_frame.width, table1._graphic_frame.height
        x2, y2, cx2, cy2 = table2._graphic_frame.left, table2._graphic_frame.top, table2._graphic_frame.width, table2._graphic_frame.height

        x = min(x1, x2)
        y = min(y1, y2)
        cx = max(x1 + cx1, x2 + cx2) - x
        cy = max(y1 + cy1, y2 + cy2) - y

        border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
        border_shape.fill.background()
        border_shape.line.color.rgb = outline_color
        border_shape.line.width = Pt(border_width)

    add_combined_border(slide, table1, table1_illum, 2, RGBColor(0, 128, 0))
    add_combined_border(slide, table2, table2_illum, 2, RGBColor(255, 215, 0))
    
    def add_red_box_over_table(slide, table, start_day, num_days):
        rows = table.rows
        cols = table.columns
        day_row_idx = 0
        gusts_row_idx = len(rows) - 1
        num_actual_cols = len(cols) - 1  # Exclude the header column

        # Ensure num_days is an integer and adjust to fit the table if needed
        num_days = int(num_days) * 2
        start_day = int(start_day) * 2 - 1  # Adjust to zero-based index
        
        if num_days + start_day > num_actual_cols:
            num_days = num_actual_cols - start_day

        # Calculate the left position and width
        x = table._graphic_frame.left + sum([table.columns[i + 1].width for i in range(start_day)])
        cx = sum([table.columns[i + 1].width for i in range(start_day, start_day + num_days)])

        # Calculate the top position and height
        y = table._graphic_frame.top
        cy = sum([table.rows[i].height for i in range(day_row_idx, gusts_row_idx + 1)])

        border_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
        border_shape.fill.background()
        border_shape.line.color.rgb = RGBColor(255, 0, 0)
        border_shape.line.width = Pt(1)


    # Add red box over the table depending on operation length
    add_red_box_over_table(slide, table1, start_day, num_days)
    add_red_box_over_table(slide, table2, start_day, num_days)

    prs.save(".output/Weather_Forecast.pptx")
