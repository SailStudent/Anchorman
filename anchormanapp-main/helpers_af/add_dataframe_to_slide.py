from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from helpers_af.merge_day_cells import merge_day_cells


def add_dataframe_to_slide_single_city(location_forecast_df, location_illum_df, slide, t_left=Inches(0), t_top= Inches(3.8), t_width = Inches(6.64), t_height= Inches(2.3), dual_city=False):
    rows, cols = location_forecast_df.shape
    left = t_left
    top = t_top
    width = t_width
    height = t_height

    table1 = slide.shapes.add_table(rows, cols + 1, left, top, width, height).table

    day_values = location_forecast_df.loc['Day'].values
    time_values = location_forecast_df.loc['Time'].values

    for i, day in enumerate(day_values):
        cell = table1.cell(0, i)
        cell.text = str(day) if day is not None else ""
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    for j, idx in enumerate(location_forecast_df.index):

        row_header_cell = table1.cell(j, 0)
        row_header_cell.text = "" if idx == 'Time' else str(idx)
        row_header_cell.text_frame.paragraphs[0].font.bold = True
        row_header_cell.text_frame.paragraphs[0].font.size = Pt(8)
        for i, col in enumerate(location_forecast_df.columns):
            table1.cell(j, i + 1).text = str(location_forecast_df.at[idx, col])
    table1.cell(2, 0).text = ""
    for cell in table1.iter_cells():
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.alignment = PP_ALIGN.CENTER
        cell.margin_left = Inches(0.00)
        cell.margin_right = Inches(0.00)
        cell.margin_top = Inches(0.01)
        cell.margin_bottom = Inches(0.01)
        cell.height = Inches(0.2)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    for col in range(1, cols + 1):
        p = table1.cell(1, col).text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(8)

    merge_day_cells(table1, day_values, time_values)

    rows, cols = location_illum_df.shape
    if dual_city:
        left = t_left
        top = t_top
        width = t_width
        height = t_height
    else:
        left = Inches(0)
        top = Inches(6.1)
        width = Inches(6.64)
        height = Inches(1.4)

    table1_illum = slide.shapes.add_table(rows + 1, cols, left, top, width, height).table

    for i, col_name in enumerate(location_illum_df.columns):
        cell = table1_illum.cell(0, i)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(8)
        
    for i, row in location_illum_df.iterrows():
        for j, value in enumerate(row):
            cell = table1_illum.cell(i + 1, j)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)

    for cell in table1_illum.iter_cells():
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(8)
        p.alignment = PP_ALIGN.CENTER
        cell.margin_left = Inches(0.00)
        cell.margin_right = Inches(0.00)
        cell.margin_top = Inches(0.01)
        cell.margin_bottom = Inches(0.01)
        cell.height = Inches(0.2)
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    
    return table1, table1_illum
